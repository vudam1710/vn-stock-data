"""
TechWorld Revenue Analysis - PPTX Slide Deck Generator
16-slide deck using python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import copy

# ─────────────────────────────────────────────
# DESIGN TOKENS
# ─────────────────────────────────────────────
BG_DARK  = RGBColor(8,  15,  30)
BG_LIGHT = RGBColor(240, 242, 248)

BLUE_PRIMARY  = RGBColor(26,  63, 204)
BLUE_BRIGHT   = RGBColor(43,  78, 255)
PURPLE        = RGBColor(124, 58, 237)
ORANGE        = RGBColor(255, 107, 43)
TEAL          = RGBColor(13,  148, 136)
AMBER         = RGBColor(245, 158, 11)
WHITE         = RGBColor(255, 255, 255)
DARK_TEXT     = RGBColor(10,  20,  50)
MID_GREY      = RGBColor(100, 110, 140)
LIGHT_GREY    = RGBColor(200, 205, 220)
RED_ALERT     = RGBColor(239,  68,  68)
GREEN_OK      = RGBColor(16, 185, 129)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height,
                text, font_size=14, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, color, thickness=2):
    """Horizontal accent line."""
    from pptx.util import Pt as PtUtil
    shape = slide.shapes.add_shape(
        1, left, top, width, Pt(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def section_label(slide, section, topic, is_dark, x=Inches(0.5), y=Inches(0.18)):
    """Small breadcrumb label top-left."""
    color = RGBColor(130, 150, 200) if is_dark else MID_GREY
    text = f"{section}  ›  {topic}" if topic else section
    add_textbox(slide, x, y, Inches(6), Inches(0.3),
                text, font_size=9, color=color)


def kpi_box(slide, left, top, width, height,
            label, value, label_color, value_color, bg_color, border_color=None):
    """Rounded-corner KPI card (simulated with rect)."""
    add_rect(slide, left, top, width, height, bg_color,
             line_color=border_color or bg_color, line_width=1)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.08),
                width - Inches(0.2), Inches(0.25),
                label, font_size=9, bold=True, color=label_color,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.05), top + Inches(0.33),
                width - Inches(0.1), height - Inches(0.4),
                value, font_size=14, bold=True, color=value_color,
                align=PP_ALIGN.CENTER)


def divider_line(slide, y, is_dark):
    color = RGBColor(40, 60, 100) if is_dark else LIGHT_GREY
    add_line(slide, Inches(0.5), y, Inches(12.33), color, thickness=1)


# ─────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────

def slide_01_cover(prs):
    """Slide 1 — Cover (dark)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_DARK)

    # Full-width top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)

    # Main title
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.7),
                "TechWorld Revenue Analysis",
                font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Hero opening hook
    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(11.5), Inches(0.55),
                "TechWorld lost $1.5M in revenue while margins held perfectly intact",
                font_size=19, bold=False, color=RGBColor(160, 180, 230), align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, Inches(0.7), Inches(1.85), Inches(11.0), Inches(0.45),
                "Why Revenue Fell -28% and What Comes Next",
                font_size=15, bold=False, color=RGBColor(120, 145, 200), align=PP_ALIGN.LEFT)

    # Hero numbers row
    add_textbox(slide, Inches(0.7), Inches(2.45), Inches(3.0), Inches(0.55),
                "-28.4%", font_size=48, bold=True, color=RED_ALERT, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.7), Inches(2.98), Inches(4.0), Inches(0.35),
                "YoY Revenue Decline", font_size=12, color=RGBColor(150, 165, 210))
    add_textbox(slide, Inches(4.8), Inches(2.45), Inches(4.0), Inches(0.55),
                "33.2%", font_size=48, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(4.8), Inches(2.98), Inches(4.0), Inches(0.35),
                "Profit Margin — Unchanged", font_size=12, color=RGBColor(150, 165, 210))

    # Horizontal divider
    add_line(slide, Inches(0.7), Inches(3.45), Inches(11.6), RGBColor(40, 60, 100), thickness=1)

    # 4 KPI boxes
    kpi_data = [
        ("PERIOD",            "Jan 2022 – Sep 2023"),
        ("TOTAL REVENUE",     "$8.75M"),
        ("COMPLETED ORDERS",  "13,925"),
        ("PROFIT MARGIN",     "33.2%"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(0.95)
    gap   = Inches(0.18)
    start_x = Inches(0.7)
    y_pos = Inches(3.6)
    for i, (lbl, val) in enumerate(kpi_data):
        lx = start_x + i * (box_w + gap)
        bg = RGBColor(18, 31, 53)
        add_rect(slide, lx, y_pos, box_w, box_h, bg,
                 line_color=RGBColor(40, 65, 120), line_width=1)
        add_textbox(slide, lx + Inches(0.1), y_pos + Inches(0.08),
                    box_w - Inches(0.2), Inches(0.25),
                    lbl, font_size=9, bold=True,
                    color=RGBColor(100, 130, 200), align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.05), y_pos + Inches(0.35),
                    box_w - Inches(0.1), Inches(0.5),
                    val, font_size=16, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)

    # Data annotations box
    add_textbox(slide, Inches(0.7), Inches(4.75), Inches(11.6), Inches(0.35),
                "Jan-Sep 2022: $5.10M  |  Jan-Sep 2023: $3.65M  |  Delta: -$1.49M  |  Orders -27.5%  |  Net Profit $2.90M",
                font_size=10, color=RGBColor(80, 105, 165), align=PP_ALIGN.CENTER)

    # Insight callout
    add_rect(slide, Inches(0.7), Inches(5.2), Inches(11.6), Inches(0.75),
             RGBColor(15, 25, 50), line_color=BLUE_PRIMARY, line_width=1)
    add_textbox(slide, Inches(0.9), Inches(5.3), Inches(11.2), Inches(0.55),
                "KEY INSIGHT  Revenue and profit fell -28% in lockstep while margin held at exactly 33.2% — this is a pure volume problem, not a pricing or cost failure.",
                font_size=11, color=RGBColor(180, 200, 240))

    # Bottom accent line
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)

    add_notes(slide, "Open by framing the paradox: revenue is down sharply but margin hasn't moved a single decimal point — 33.2% both years. That tells the audience immediately this is not a pricing or cost problem, it is a pure volume problem. Every dollar TechWorld made it kept 33 cents of; they are just making fewer dollars.")
    return slide


def slide_02_background(prs):
    """Slide 2 — Background (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    section_label(slide, "Context", "Business Context & Analysis Scope", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.55),
                "21 months of order data across 4 regions and categories",
                font_size=22, bold=True, color=DARK_TEXT)

    add_textbox(slide, Inches(0.5), Inches(0.95), Inches(12.2), Inches(0.4),
                "TechWorld is a focused tech e-commerce business with one dominant risk",
                font_size=13, color=MID_GREY)

    divider_line(slide, Inches(1.4), False)

    # 6 stat boxes — 3 per row
    stats = [
        ("20,429", "Raw Dataset Rows",      BLUE_PRIMARY),
        ("13,925", "Completed Orders",       TEAL),
        ("4",      "Regions",               PURPLE),
        ("4",      "Product Categories",    AMBER),
        ("5",      "Suppliers",             ORANGE),
        ("B  82/100", "Data Quality Grade", GREEN_OK),
    ]
    box_w = Inches(3.8)
    box_h = Inches(1.3)
    gap_x = Inches(0.28)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(1.6)
    for i, (val, lbl, accent) in enumerate(stats):
        col = i % 3
        row = i // 3
        lx = start_x + col * (box_w + gap_x)
        ly = start_y + row * (box_h + gap_y)
        bg = RGBColor(228, 232, 248)
        add_rect(slide, lx, ly, box_w, box_h, bg,
                 line_color=RGBColor(190, 200, 230), line_width=1)
        # accent left bar
        add_rect(slide, lx, ly, Inches(0.06), box_h, accent)
        add_textbox(slide, lx + Inches(0.18), ly + Inches(0.12),
                    box_w - Inches(0.3), Inches(0.65),
                    val, font_size=28, bold=True, color=accent)
        add_textbox(slide, lx + Inches(0.18), ly + Inches(0.78),
                    box_w - Inches(0.3), Inches(0.4),
                    lbl, font_size=11, color=MID_GREY)

    # Right panel text
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.2), Inches(0.35),
                "Analysis scope: 3 analytical layers — Descriptive (what happened) · Diagnostic (why) · Forecasting (what comes next)",
                font_size=11, color=MID_GREY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.45), Inches(12.2), Inches(0.35),
                "Data filters: 15 rows with invalid 2026 dates excluded · Order ID 20422 corrupted · Return rate 2.8% confirms revenue integrity",
                font_size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    add_notes(slide, "Set the scene cleanly: this is a transactional order-level dataset, fully cleaned and filtered. The three-layer analysis structure tells the audience what level of rigour they are about to see — descriptive facts, causal diagnosis, and forward-looking forecast. Emphasise the 2.8% return rate: this is a high-quality transactional dataset where revenue figures mean what they say.")
    return slide


def slide_03_framework(prs):
    """Slide 3 — Framework (dark)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    section_label(slide, "Approach", "Analysis Framework", True)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.55),
                "Three analytical layers converge on one diagnosis and one forecast",
                font_size=24, bold=True, color=WHITE)

    divider_line(slide, Inches(1.0), True)

    cols = [
        {
            "num": "01", "title": "What Happened", "accent": TEAL,
            "icon": "📊",
            "bullets": [
                "Revenue fell -28.4% YoY to $3.65M",
                "Profit margin held flat at 33.2% — pure volume collapse",
                "Monthly slope: -$1,081/month across 21 months",
            ],
            "footer": "Descriptive layer: trend, segment, region",
            "stat": "$5.10M → $3.65M",
        },
        {
            "num": "02", "title": "Why It Happened", "accent": ORANGE,
            "icon": "🔍",
            "bullets": [
                "Volume drives 70% of revenue variance (r=0.818)",
                "87.3% Electronics concentration — no category buffer",
                "East -35.3% YoY / South Wearables -58.6%",
            ],
            "footer": "Diagnostic layer: 5 hypotheses, 3 rejected",
            "stat": "70% variance explained",
        },
        {
            "num": "03", "title": "What Comes Next", "accent": BLUE_BRIGHT,
            "icon": "→",
            "bullets": [
                "Q4 2023 ensemble forecast: $1.22M total",
                "3 models agree within 0.3 MAPE points",
                "No reversal signal — intervention required",
            ],
            "footer": "Forecasting: 3-model ensemble, MAPE <4%",
            "stat": "$1.22M Q4 forecast",
        },
    ]

    card_w = Inches(3.9)
    card_h = Inches(5.5)
    gap    = Inches(0.22)
    start_x = Inches(0.55)
    y0 = Inches(1.1)

    for i, col in enumerate(cols):
        lx = start_x + i * (card_w + gap)
        bg = RGBColor(14, 23, 45)
        add_rect(slide, lx, y0, card_w, card_h, bg,
                 line_color=col["accent"], line_width=1)
        # Top accent strip
        add_rect(slide, lx, y0, card_w, Inches(0.07), col["accent"])
        # Number badge
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.18),
                    Inches(0.6), Inches(0.4),
                    col["num"], font_size=22, bold=True, color=col["accent"])
        # Title
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.6),
                    card_w - Inches(0.3), Inches(0.45),
                    col["title"], font_size=17, bold=True, color=WHITE)
        # Stat
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(1.08),
                    card_w - Inches(0.3), Inches(0.38),
                    col["stat"], font_size=12, bold=True, color=col["accent"])
        # Bullets
        for j, bullet in enumerate(col["bullets"]):
            add_textbox(slide, lx + Inches(0.22), y0 + Inches(1.55) + j * Inches(0.78),
                        card_w - Inches(0.4), Inches(0.7),
                        f"• {bullet}", font_size=11, color=RGBColor(180, 195, 230))
        # Footer
        add_rect(slide, lx, y0 + card_h - Inches(0.45), card_w, Inches(0.45),
                 RGBColor(10, 18, 38))
        add_textbox(slide, lx + Inches(0.15), y0 + card_h - Inches(0.38),
                    card_w - Inches(0.25), Inches(0.32),
                    col["footer"], font_size=9, color=MID_GREY)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    add_notes(slide, "This is the map before the journey. Walk through all three layers at the structural level before diving into detail. The purpose is to set expectations: every finding in this deck traces back to one of these three questions, and the conclusion from each layer will align.")
    return slide


def slide_04_yoy(prs):
    """Slide 4 — YoY Comparison (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    section_label(slide, "Context", "What's Happening", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Revenue fell -28% YoY while margin held — confirming a volume collapse not a pricing failure",
                font_size=20, bold=True, color=DARK_TEXT)

    divider_line(slide, Inches(0.95), False)

    # Hero KPIs
    heroes = [
        ("-28.4%",  "YoY Revenue Decline",    RED_ALERT),
        ("33.2%",   "Margin — Unchanged",      TEAL),
        ("-$1.49M", "Revenue Gap Jan-Sep",     ORANGE),
        ("-27.5%",  "Orders Decline YoY",      PURPLE),
    ]
    hw = Inches(2.9)
    hh = Inches(0.9)
    hgap = Inches(0.25)
    hx = Inches(0.5)
    hy = Inches(1.1)
    for i, (val, lbl, clr) in enumerate(heroes):
        lx = hx + i * (hw + hgap)
        add_rect(slide, lx, hy, hw, hh, RGBColor(235, 238, 252),
                 line_color=clr, line_width=2)
        add_textbox(slide, lx + Inches(0.1), hy + Inches(0.04),
                    hw - Inches(0.2), Inches(0.5),
                    val, font_size=26, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.1), hy + Inches(0.54),
                    hw - Inches(0.2), Inches(0.28),
                    lbl, font_size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

    divider_line(slide, Inches(2.12), False)

    # YoY comparison table
    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(6.0), Inches(0.35),
                "Year-over-Year Comparison (Jan–Sep)", font_size=13, bold=True, color=DARK_TEXT)

    headers = ["Metric", "Jan-Sep 2022", "Jan-Sep 2023", "Change"]
    rows = [
        ["Revenue",       "$5,101,180",  "$3,650,160",  "-28.4%  ▼"],
        ["Net Profit",    "$1,693,520",  "$1,211,203",  "-28.5%  ▼"],
        ["Orders",        "8,073",       "5,852",        "-27.5%  ▼"],
        ["Units Sold",    "~13,443",     "~9,714",       "-27.7%  ▼"],
        ["Profit Margin", "33.20%",      "33.18%",       "FLAT  ●"],
    ]

    col_widths = [Inches(2.0), Inches(1.6), Inches(1.6), Inches(1.4)]
    col_starts = [Inches(0.5), Inches(2.55), Inches(4.2), Inches(5.85)]
    row_h = Inches(0.38)
    header_y = Inches(2.62)

    # Header row bg
    add_rect(slide, Inches(0.5), header_y, Inches(6.8), row_h, BLUE_PRIMARY)
    for j, (hdr, cw, cx) in enumerate(zip(headers, col_widths, col_starts)):
        add_textbox(slide, cx + Inches(0.05), header_y + Inches(0.07),
                    cw - Inches(0.1), row_h - Inches(0.1),
                    hdr, font_size=10, bold=True, color=WHITE)

    for i, row in enumerate(rows):
        ry = header_y + (i + 1) * row_h
        bg = RGBColor(228, 232, 248) if i % 2 == 0 else RGBColor(240, 242, 248)
        add_rect(slide, Inches(0.5), ry, Inches(6.8), row_h, bg)
        for j, (cell, cw, cx) in enumerate(zip(row, col_widths, col_starts)):
            color = DARK_TEXT
            if j == 3:
                if "▼" in cell:
                    color = RED_ALERT
                elif "FLAT" in cell:
                    color = TEAL
            add_textbox(slide, cx + Inches(0.05), ry + Inches(0.07),
                        cw - Inches(0.1), row_h - Inches(0.1),
                        cell, font_size=10, bold=(j == 0), color=color)

    # Right panel — insight text
    add_rect(slide, Inches(7.6), Inches(2.2), Inches(5.2), Inches(4.0),
             RGBColor(228, 232, 248), line_color=BLUE_PRIMARY, line_width=1)
    add_textbox(slide, Inches(7.8), Inches(2.3), Inches(4.8), Inches(0.35),
                "The Volume Fingerprint", font_size=14, bold=True, color=DARK_TEXT)
    insights = [
        "All 3 headline metrics fell 27-29% YoY.",
        "This proportional decline is the signature",
        "of pure volume loss — not pricing failure.",
        "",
        "Profit margin: 33.20% → 33.18%",
        "A 0.02pp change is not a rounding diff —",
        "it is proof that pricing held perfectly.",
        "",
        "Monthly revenue slope: -$1,081/month",
        "Peak: Jan 2022 $468K | Trough: Apr 2023 $365K",
        "CV = 6.5% — moderate, not extreme volatility",
    ]
    for k, line in enumerate(insights):
        clr = BLUE_PRIMARY if line.startswith("Profit") or line.startswith("Monthly") else (DARK_TEXT if line else DARK_TEXT)
        bold = line.startswith("Profit") or line.startswith("Monthly")
        add_textbox(slide, Inches(7.8), Inches(2.75) + k * Inches(0.29),
                    Inches(4.8), Inches(0.28),
                    line, font_size=10, bold=bold,
                    color=BLUE_PRIMARY if bold else DARK_TEXT)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    add_notes(slide, "This is the single most important diagnostic slide in the deck. Spend time on the margin number — 33.20% versus 33.18% is not a rounding difference, it is a proof. If TechWorld had cut prices to stimulate volume, margin would have compressed. Every dollar of revenue TechWorld earned in 2023 was just as profitable as in 2022 — there were simply fewer of those dollars.")
    return slide


def slide_05_category_region(prs):
    """Slide 5 — Category + Region (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    section_label(slide, "Context", "What's Happening", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Electronics at 87% of revenue means one category IS the entire business",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # LEFT — Category table
    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(5.8), Inches(0.35),
                "Revenue by Category", font_size=13, bold=True, color=DARK_TEXT)

    cat_headers = ["Category", "Revenue", "Share", "Margin", "YoY"]
    cat_rows = [
        ["Electronics",  "$7,630,300", "87.2%", "32.4%", "-28.3%"],
        ["Wearables",    "$424,800",   "4.9%",  "37.9%", "-32.2%"],
        ["Audio",        "$411,200",   "4.7%",  "38.3%", "-27.4%"],
        ["Accessories",  "$285,040",   "3.3%",  "40.7%", "-26.7%"],
    ]
    cat_col_w = [Inches(1.5), Inches(1.3), Inches(0.75), Inches(0.75), Inches(0.85)]
    cat_col_x = [Inches(0.5), Inches(2.05), Inches(3.4), Inches(4.18), Inches(4.96)]
    row_h = Inches(0.4)
    hy = Inches(1.48)

    add_rect(slide, Inches(0.5), hy, Inches(5.6), row_h, BLUE_PRIMARY)
    for j, (h, cw, cx) in enumerate(zip(cat_headers, cat_col_w, cat_col_x)):
        add_textbox(slide, cx + Inches(0.04), hy + Inches(0.07),
                    cw - Inches(0.08), row_h - Inches(0.1),
                    h, font_size=9, bold=True, color=WHITE)

    for i, row in enumerate(cat_rows):
        ry = hy + (i + 1) * row_h
        bg = RGBColor(220, 225, 245) if i == 0 else (RGBColor(228, 232, 248) if i % 2 == 0 else BG_LIGHT)
        add_rect(slide, Inches(0.5), ry, Inches(5.6), row_h, bg)
        # highlight bar for Electronics
        if i == 0:
            add_rect(slide, Inches(0.5), ry, Inches(0.06), row_h, BLUE_BRIGHT)
        for j, (cell, cw, cx) in enumerate(zip(row, cat_col_w, cat_col_x)):
            clr = BLUE_BRIGHT if (i == 0 and j == 0) else (RED_ALERT if (j == 4 and "%" in cell) else DARK_TEXT)
            add_textbox(slide, cx + Inches(0.04), ry + Inches(0.07),
                        cw - Inches(0.08), row_h - Inches(0.1),
                        cell, font_size=9, bold=(i == 0), color=clr)

    # Callout box
    add_rect(slide, Inches(0.5), Inches(3.2), Inches(5.6), Inches(0.7),
             RGBColor(230, 235, 255), line_color=BLUE_BRIGHT, line_width=1)
    add_textbox(slide, Inches(0.65), Inches(3.28), Inches(5.3), Inches(0.55),
                "87.2% concentration — the smallest categories carry HIGHER margins (38-41%) but are too small to buffer any Electronics shock",
                font_size=10, color=BLUE_PRIMARY)

    # RIGHT — Region table
    add_textbox(slide, Inches(6.9), Inches(1.05), Inches(5.8), Inches(0.35),
                "Revenue by Region (with growth slope)", font_size=13, bold=True, color=DARK_TEXT)

    reg_headers = ["Region", "Revenue", "Share", "Slope/mo", "YoY"]
    reg_rows = [
        ["East",  "$2,226,080", "25.4%", "-$1,114", "-35.3%"],
        ["West",  "$2,202,350", "25.2%", "+$609",   "-22.3%"],
        ["North", "$2,169,460", "24.8%", "-$192",   "-27.7%"],
        ["South", "$2,153,450", "24.6%", "-$384",   "-27.8%"],
    ]
    reg_col_w = [Inches(1.3), Inches(1.4), Inches(0.8), Inches(0.95), Inches(0.85)]
    reg_col_x = [Inches(6.9), Inches(8.25), Inches(9.7), Inches(10.53), Inches(11.52)]
    hy2 = Inches(1.48)

    add_rect(slide, Inches(6.9), hy2, Inches(5.6), row_h, BLUE_PRIMARY)
    for j, (h, cw, cx) in enumerate(zip(reg_headers, reg_col_w, reg_col_x)):
        add_textbox(slide, cx + Inches(0.04), hy2 + Inches(0.07),
                    cw - Inches(0.08), row_h - Inches(0.1),
                    h, font_size=9, bold=True, color=WHITE)

    slope_colors = [RED_ALERT, TEAL, MID_GREY, ORANGE]
    for i, row in enumerate(reg_rows):
        ry = hy2 + (i + 1) * row_h
        bg = RGBColor(228, 232, 248) if i % 2 == 0 else BG_LIGHT
        add_rect(slide, Inches(6.9), ry, Inches(5.6), row_h, bg)
        for j, (cell, cw, cx) in enumerate(zip(row, reg_col_w, reg_col_x)):
            clr = DARK_TEXT
            if j == 3:
                clr = slope_colors[i]
            elif j == 4:
                clr = RED_ALERT if "-" in cell else TEAL
            add_textbox(slide, cx + Inches(0.04), ry + Inches(0.07),
                        cw - Inches(0.08), row_h - Inches(0.1),
                        cell, font_size=9, bold=(j == 3), color=clr)

    # West callout
    add_rect(slide, Inches(6.9), Inches(3.2), Inches(5.6), Inches(0.7),
             RGBColor(220, 248, 240), line_color=TEAL, line_width=1)
    add_textbox(slide, Inches(7.05), Inches(3.28), Inches(5.3), Inches(0.55),
                "West (+$609/month) is the ONLY growing region — proof that growth is possible. East at -$1,114/month is nearly 2x worse than the aggregate.",
                font_size=10, color=TEAL)

    # Visual bar chart simulation for categories
    add_textbox(slide, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.3),
                "Electronics share of revenue across all 21 months:", font_size=11, bold=True, color=DARK_TEXT)
    bar_y = Inches(4.4)
    bar_h = Inches(0.45)
    total_w = Inches(12.0)
    # Electronics bar
    elec_w = total_w * 0.872
    add_rect(slide, Inches(0.5), bar_y, elec_w, bar_h, BLUE_BRIGHT)
    add_textbox(slide, Inches(0.55), bar_y + Inches(0.08), elec_w - Inches(0.2), bar_h - Inches(0.1),
                "Electronics  87.2%", font_size=11, bold=True, color=WHITE)
    others = [
        ("Wearables 4.9%", 0.049, PURPLE),
        ("Audio 4.7%", 0.047, TEAL),
        ("Acc 3.3%", 0.033, AMBER),
    ]
    cur_x = Inches(0.5) + elec_w
    for lbl, frac, clr in others:
        w = total_w * frac
        add_rect(slide, cur_x, bar_y, w, bar_h, clr)
        if w > Inches(0.5):
            add_textbox(slide, cur_x + Inches(0.02), bar_y + Inches(0.1), w - Inches(0.04), bar_h - Inches(0.15),
                        lbl, font_size=7, color=WHITE)
        cur_x += w

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    add_notes(slide, "Two structural facts the management team must internalise before the diagnostic section. First: the business has no diversification buffer — 87% in one category is not a product mix. Second: the aggregate -28% figure conceals that one region is actually growing. West at +$609/month is proof that TechWorld can generate growth in the right conditions.")
    return slide


def slide_06_segment(prs):
    """Slide 6 — Segment & Seasonality (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    section_label(slide, "Context", "What's Happening", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "14% of orders generate 44% of revenue — the high-value segment is the business engine",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # LEFT — segment table
    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(5.8), Inches(0.35),
                "Order Segment Analysis", font_size=13, bold=True, color=DARK_TEXT)

    seg_headers = ["Segment", "Orders", "Share", "Revenue", "AOV", "Margin"]
    seg_rows = [
        ["High-Value (>$1K)", "1,950", "14.0%", "$3,869,000", "$1,984", "32.7%"],
        ["Standard (<$1K)",   "11,975", "86.0%", "$4,882,340", "$408",   "33.6%"],
    ]
    seg_cw = [Inches(2.0), Inches(0.7), Inches(0.7), Inches(1.4), Inches(0.8), Inches(0.75)]
    seg_cx = [Inches(0.5), Inches(2.55), Inches(3.28), Inches(4.02), Inches(5.46), Inches(6.3)]
    row_h = Inches(0.45)
    hy = Inches(1.48)
    add_rect(slide, Inches(0.5), hy, Inches(6.6), row_h, BLUE_PRIMARY)
    for j, (h, cw, cx) in enumerate(zip(seg_headers, seg_cw, seg_cx)):
        add_textbox(slide, cx + Inches(0.03), hy + Inches(0.08),
                    cw - Inches(0.06), row_h - Inches(0.1),
                    h, font_size=8, bold=True, color=WHITE)
    for i, row in enumerate(seg_rows):
        ry = hy + (i + 1) * row_h
        bg = RGBColor(220, 225, 245) if i == 0 else RGBColor(236, 240, 252)
        add_rect(slide, Inches(0.5), ry, Inches(6.6), row_h, bg)
        if i == 0:
            add_rect(slide, Inches(0.5), ry, Inches(0.06), row_h, AMBER)
        for j, (cell, cw, cx) in enumerate(zip(row, seg_cw, seg_cx)):
            clr = AMBER if (i == 0 and j in [1, 3, 4]) else DARK_TEXT
            add_textbox(slide, cx + Inches(0.03), ry + Inches(0.08),
                        cw - Inches(0.06), row_h - Inches(0.1),
                        cell, font_size=9, bold=(j == 0), color=clr)

    # Leverage insight
    add_rect(slide, Inches(0.5), Inches(2.48), Inches(6.6), Inches(0.75),
             RGBColor(255, 248, 230), line_color=AMBER, line_width=1)
    add_textbox(slide, Inches(0.65), Inches(2.55), Inches(6.3), Inches(0.6),
                "LEVERAGE  Losing 200 high-value orders equals losing 975 standard orders in revenue impact. Protecting this segment is the #1 retention priority.",
                font_size=10, color=RGBColor(140, 90, 10))

    # Mini bar for segment
    add_textbox(slide, Inches(0.5), Inches(3.38), Inches(6.6), Inches(0.28),
                "Revenue share by segment:", font_size=10, bold=True, color=DARK_TEXT)
    bar_y = Inches(3.7)
    bar_h = Inches(0.35)
    add_rect(slide, Inches(0.5), bar_y, Inches(6.6 * 0.442), bar_h, AMBER)
    add_textbox(slide, Inches(0.55), bar_y + Inches(0.07), Inches(2.5), Inches(0.25),
                "High-Value 44%", font_size=9, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5) + Inches(6.6 * 0.442), bar_y, Inches(6.6 * 0.558), bar_h, RGBColor(160, 175, 220))
    add_textbox(slide, Inches(0.5) + Inches(6.6 * 0.442) + Inches(0.05), bar_y + Inches(0.07), Inches(2.5), Inches(0.25),
                "Standard 56%", font_size=9, color=WHITE)

    # RIGHT — February seasonality
    add_textbox(slide, Inches(7.4), Inches(1.05), Inches(5.5), Inches(0.35),
                "February Seasonality Pattern", font_size=13, bold=True, color=DARK_TEXT)

    feb_data = [
        ("Feb 2022 MoM Drop",     "-21.4%",  RED_ALERT),
        ("Calendar Effect",        "-9.7%",   ORANGE),
        ("Genuine Demand Softness","-11.7%",  RED_ALERT),
        ("North Feb 2022 MoM",     "-36.9%",  RED_ALERT),
        ("Feb 2023 MoM Drop",      "-3.2%",   ORANGE),
    ]
    for i, (lbl, val, clr) in enumerate(feb_data):
        ry = Inches(1.5) + i * Inches(0.52)
        bg = RGBColor(255, 240, 238) if "Genuine" in lbl or "Feb 2022" in lbl else RGBColor(235, 238, 252)
        add_rect(slide, Inches(7.4), ry, Inches(5.5), Inches(0.46), bg,
                 line_color=clr, line_width=1)
        add_textbox(slide, Inches(7.55), ry + Inches(0.06),
                    Inches(3.5), Inches(0.32), lbl, font_size=10, color=DARK_TEXT)
        add_textbox(slide, Inches(11.2), ry + Inches(0.04),
                    Inches(1.5), Inches(0.38), val, font_size=16, bold=True,
                    color=clr, align=PP_ALIGN.RIGHT)

    add_rect(slide, Inches(7.4), Inches(4.2), Inches(5.5), Inches(0.9),
             RGBColor(255, 240, 238), line_color=RED_ALERT, line_width=1)
    add_textbox(slide, Inches(7.55), Inches(4.28), Inches(5.2), Inches(0.75),
                "PLANNABLE EVENT  The Feb-to-Mar pattern repeated in both 2022 and 2023. A February promotion for North region can buffer this structural annual trough.",
                font_size=10, color=RED_ALERT)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    add_notes(slide, "Two insights with very different action implications. The segment finding tells management that protecting or growing 200 high-value accounts has 5x the revenue impact of equivalent standard order growth. The February finding is about planning: this is a calendar-driven structural trough that repeats every year. A promotions campaign designed before January would buffer the inevitable February softness.")
    return slide


def slide_07_findings_overview(prs):
    """Slide 7 — Findings Overview (dark)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), RED_ALERT)
    section_label(slide, "Findings", "Overview", True)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Three root causes explain the revenue collapse — all confirmed at high confidence",
                font_size=24, bold=True, color=WHITE)
    divider_line(slide, Inches(1.0), True)

    causes = [
        {
            "num": "01", "title": "Demand Volume Drop",
            "accent": RED_ALERT,
            "stat": "70% of variance",
            "evidence": "r=0.818 correlation",
            "bullets": [
                "Order volume drives 70% of revenue variance",
                "Quantity-Revenue correlation: r=0.818",
                "All 3 volatile months driven by 15-22% order swings",
                "Price is stable ($479-$539 range) — not the lever",
            ],
            "footer": "H1 CONFIRMED — HIGH confidence",
        },
        {
            "num": "02", "title": "Electronics Concentration",
            "accent": ORANGE,
            "stat": "87.3% share",
            "evidence": "Frozen 30 months",
            "bullets": [
                "Electronics = 87.3% of revenue",
                "Category mix frozen: std dev only 0.9 ppts",
                "Every Electronics shock hits total revenue 1:1",
                "Feb 2022: Electronics -21.8% vs total -21.4%",
            ],
            "footer": "H3 CONFIRMED — HIGH confidence",
        },
        {
            "num": "03", "title": "Seasonal & Regional Gaps",
            "accent": AMBER,
            "stat": "-58.6% Wearables South",
            "evidence": "East -35.3% YoY",
            "bullets": [
                "Feb trough: -21.4% in 2022, -3.2% in 2023",
                "South Wearables collapsed -58.6%",
                "East -35.3% YoY — worst region",
                "West +609/month — only growth signal",
            ],
            "footer": "H2+H4 CONFIRMED — MEDIUM-HIGH",
        },
    ]

    card_w = Inches(3.9)
    card_h = Inches(5.6)
    gap    = Inches(0.22)
    start_x = Inches(0.55)
    y0 = Inches(1.1)

    for i, c in enumerate(causes):
        lx = start_x + i * (card_w + gap)
        add_rect(slide, lx, y0, card_w, card_h, RGBColor(14, 23, 45),
                 line_color=c["accent"], line_width=1)
        add_rect(slide, lx, y0, card_w, Inches(0.07), c["accent"])
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.15),
                    Inches(0.55), Inches(0.38),
                    c["num"], font_size=20, bold=True, color=c["accent"])
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.58),
                    card_w - Inches(0.3), Inches(0.42),
                    c["title"], font_size=16, bold=True, color=WHITE)
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(1.05),
                    card_w - Inches(0.3), Inches(0.33),
                    c["stat"], font_size=13, bold=True, color=c["accent"])
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(1.38),
                    card_w - Inches(0.3), Inches(0.28),
                    c["evidence"], font_size=10, color=MID_GREY)
        for j, bullet in enumerate(c["bullets"]):
            add_textbox(slide, lx + Inches(0.22), y0 + Inches(1.78) + j * Inches(0.72),
                        card_w - Inches(0.4), Inches(0.66),
                        f"• {bullet}", font_size=11, color=RGBColor(180, 195, 230))
        add_rect(slide, lx, y0 + card_h - Inches(0.45), card_w, Inches(0.45),
                 RGBColor(10, 18, 38))
        add_textbox(slide, lx + Inches(0.15), y0 + card_h - Inches(0.38),
                    card_w - Inches(0.25), Inches(0.32),
                    c["footer"], font_size=9, color=MID_GREY)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), RED_ALERT)
    add_notes(slide, "This is the pivot slide — the moment the deck moves from observing what happened to explaining why. Three root causes with very different action implications: the first is a demand generation problem; the second is a portfolio diversification problem; the third is a targeted operational fix. Each card will be expanded in detail over the next three slides.")
    return slide


def slide_08_volume_decomp(prs):
    """Slide 8 — Volume Decomposition (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), RED_ALERT)
    section_label(slide, "Root Cause #1", "Demand Volume", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Demand volume drives 70% of revenue variance — not price and not mix",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # Left — decomposition
    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(5.8), Inches(0.35),
                "Revenue Variance Decomposition (Laspeyres method)", font_size=12, bold=True, color=DARK_TEXT)

    decomp = [
        ("Volume Effect",  69.9, RED_ALERT,     "Order count & quantity changes"),
        ("Price Effect",   30.1, RGBColor(150,165,220), "Avg unit price variation ($479-$539)"),
        ("Mix Effect",     0.0,  LIGHT_GREY,     "Category share changes (< 1%)"),
    ]
    bar_x = Inches(0.5)
    max_bar_w = Inches(5.5)
    for i, (lbl, pct, clr, note) in enumerate(decomp):
        by = Inches(1.55) + i * Inches(1.1)
        bar_w = max_bar_w * (pct / 100.0)
        add_textbox(slide, bar_x, by, Inches(5.8), Inches(0.28),
                    lbl, font_size=11, bold=True, color=DARK_TEXT)
        add_rect(slide, bar_x, by + Inches(0.3), max_bar_w, Inches(0.38), RGBColor(220, 225, 240))
        if bar_w > Inches(0.1):
            add_rect(slide, bar_x, by + Inches(0.3), bar_w, Inches(0.38), clr)
        pct_label = f"{pct:.1f}%" if pct > 0 else "< 1%"
        add_textbox(slide, bar_x + max_bar_w + Inches(0.1), by + Inches(0.3),
                    Inches(0.8), Inches(0.38),
                    pct_label, font_size=13, bold=True, color=clr)
        add_textbox(slide, bar_x, by + Inches(0.72), Inches(5.8), Inches(0.28),
                    note, font_size=9, color=MID_GREY)

    # Correlation callouts
    corr_y = Inches(4.7)
    corr_data = [
        ("Quantity ↔ Revenue", "r = 0.818", RED_ALERT),
        ("Orders ↔ Revenue",   "r = 0.774", ORANGE),
        ("Price ↔ Revenue",    "r = 0.470", MID_GREY),
        ("Mkt Lag-1 ↔ Revenue","r = -0.33", PURPLE),
    ]
    for i, (lbl, val, clr) in enumerate(corr_data):
        cx = Inches(0.5) + i * Inches(1.55)
        add_rect(slide, cx, corr_y, Inches(1.45), Inches(0.7),
                 RGBColor(228, 232, 248), line_color=clr, line_width=1)
        add_textbox(slide, cx + Inches(0.05), corr_y + Inches(0.04),
                    Inches(1.35), Inches(0.32),
                    val, font_size=14, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx + Inches(0.05), corr_y + Inches(0.38),
                    Inches(1.35), Inches(0.26),
                    lbl, font_size=8, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Right — volatile months
    add_textbox(slide, Inches(7.0), Inches(1.05), Inches(5.8), Inches(0.35),
                "Volatile Months Confirm Volume Mechanism", font_size=12, bold=True, color=DARK_TEXT)

    vol_months = [
        ("Feb 2022", "-21.4% MoM", "Orders fell -18.2% → Revenue fell -21.4%", RED_ALERT),
        ("Mar 2022", "+21.0% MoM", "Orders rose +22.7% → Revenue rose +21.0%", GREEN_OK),
        ("May 2023", "+21.7% MoM", "Orders rose +13.9% → Revenue rose +21.7%", GREEN_OK),
    ]
    for i, (m, pct, note, clr) in enumerate(vol_months):
        my = Inches(1.5) + i * Inches(1.15)
        add_rect(slide, Inches(7.0), my, Inches(5.8), Inches(1.05),
                 RGBColor(228, 232, 248), line_color=clr, line_width=1)
        add_rect(slide, Inches(7.0), my, Inches(0.06), Inches(1.05), clr)
        add_textbox(slide, Inches(7.15), my + Inches(0.08),
                    Inches(2.5), Inches(0.35), m, font_size=12, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(9.8), my + Inches(0.04),
                    Inches(2.8), Inches(0.45), pct, font_size=22, bold=True,
                    color=clr, align=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(7.15), my + Inches(0.5),
                    Inches(5.5), Inches(0.42), note, font_size=10, color=MID_GREY)

    add_rect(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(0.8),
             RGBColor(255, 240, 238), line_color=RED_ALERT, line_width=1)
    add_textbox(slide, Inches(7.15), Inches(5.08), Inches(5.5), Inches(0.65),
                "KEY  Marketing spend does NOT lead revenue (lag-1 = -0.33). Increasing Google Ads in low-demand months compresses margins without generating demand.",
                font_size=10, color=RED_ALERT)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), RED_ALERT)
    add_notes(slide, "The key insight to deliver here is counterintuitive: the solution is NOT to increase marketing spend. The lag-1 correlation of -0.33 shows that spend follows revenue movements, not the other way around — this is reactive budgeting, not causal investment. The lever is demand generation at the structural level — seasonal promotions timed correctly and activation of channels with genuine demand-generation capability like Email.")
    return slide


def slide_09_concentration(prs):
    """Slide 9 — Electronics Concentration (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), RED_ALERT)
    section_label(slide, "Root Cause #2", "Concentration Risk", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "87% Electronics concentration means a single category shock collapses the entire business",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # Left — concentration facts
    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(5.9), Inches(0.35),
                "Category Concentration — 30-Month View", font_size=12, bold=True, color=DARK_TEXT)

    conc_data = [
        ("Electronics",   87.19, 32.38, BLUE_BRIGHT),
        ("Wearables",      4.85, 37.85, PURPLE),
        ("Audio",          4.70, 38.26, TEAL),
        ("Accessories",    3.26, 40.65, AMBER),
    ]
    for i, (name, share, margin, clr) in enumerate(conc_data):
        by = Inches(1.5) + i * Inches(0.9)
        bar_w = Inches(5.5) * (share / 100.0)
        bg = RGBColor(215, 220, 245) if i == 0 else RGBColor(235, 238, 252)
        add_rect(slide, Inches(0.5), by, Inches(5.5), Inches(0.52), bg)
        add_rect(slide, Inches(0.5), by, max(bar_w, Inches(0.1)), Inches(0.52), clr)
        label = f"{name}   {share:.1f}%   margin {margin:.1f}%"
        txt_color = WHITE if i == 0 and bar_w > Inches(2.0) else DARK_TEXT
        add_textbox(slide, Inches(0.55), by + Inches(0.13),
                    Inches(5.3), Inches(0.26), label, font_size=10, bold=(i==0), color=txt_color)

    # Frozen share callout
    add_rect(slide, Inches(0.5), Inches(5.2), Inches(5.9), Inches(0.8),
             RGBColor(220, 225, 248), line_color=BLUE_BRIGHT, line_width=2)
    add_textbox(slide, Inches(0.65), Inches(5.28), Inches(5.6), Inches(0.65),
                "FROZEN FOR 30 MONTHS  Electronics share range: 85.1%-88.7%. Standard deviation: only 0.9 ppts. No natural diversification is occurring.",
                font_size=10, bold=False, color=BLUE_PRIMARY)

    # Right — amplification evidence
    add_textbox(slide, Inches(7.0), Inches(1.05), Inches(5.8), Inches(0.35),
                "Amplification Effect — Volatile Months", font_size=12, bold=True, color=DARK_TEXT)

    amp_data = [
        ("Feb 2022", "Electronics", "-21.8%", "Total Revenue", "-21.4%"),
        ("Mar 2022", "Electronics", "+18.9%", "Total Revenue", "+21.0%"),
        ("May 2023", "Electronics", "+22.2%", "Total Revenue", "+21.7%"),
    ]
    for i, (month, cat1, chg1, cat2, chg2) in enumerate(amp_data):
        my = Inches(1.5) + i * Inches(1.0)
        add_rect(slide, Inches(7.0), my, Inches(5.8), Inches(0.88),
                 RGBColor(228, 232, 248), line_color=BLUE_BRIGHT, line_width=1)
        add_textbox(slide, Inches(7.1), my + Inches(0.06),
                    Inches(1.5), Inches(0.3), month, font_size=11, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(7.1), my + Inches(0.42), Inches(2.5), Inches(0.3),
                    f"{cat1}: {chg1}", font_size=10, color=BLUE_BRIGHT)
        add_textbox(slide, Inches(9.8), my + Inches(0.42), Inches(2.8), Inches(0.3),
                    f"{cat2}: {chg2}", font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)

    # Margin comparison
    add_textbox(slide, Inches(7.0), Inches(4.6), Inches(5.8), Inches(0.35),
                "Margin Premium — Grow These Categories", font_size=12, bold=True, color=DARK_TEXT)

    marg_data = [
        ("Accessories",  40.65, AMBER,      "+8.3pp vs Electronics"),
        ("Audio",        38.26, TEAL,       "+5.9pp vs Electronics"),
        ("Wearables",    37.85, PURPLE,     "+5.5pp vs Electronics"),
        ("Electronics",  32.38, MID_GREY,   "BASELINE — lowest margin"),
    ]
    for i, (name, marg, clr, note) in enumerate(marg_data):
        my = Inches(5.02) + i * Inches(0.42)
        add_rect(slide, Inches(7.0), my, Inches(5.8), Inches(0.38), RGBColor(228, 232, 248))
        bar_w = Inches(3.5) * (marg / 45.0)
        add_rect(slide, Inches(7.0), my, bar_w, Inches(0.38), clr)
        add_textbox(slide, Inches(7.05), my + Inches(0.06),
                    Inches(3.3), Inches(0.26), f"{name}  {marg:.1f}%", font_size=9,
                    bold=True, color=WHITE if marg > 35 else DARK_TEXT)
        add_textbox(slide, Inches(10.7), my + Inches(0.06),
                    Inches(2.0), Inches(0.26), note, font_size=8, color=MID_GREY)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), RED_ALERT)
    add_notes(slide, "The structural risk insight here is important to frame correctly: this is not about whether Electronics is a good or bad category. It is about portfolio resilience. The fix is not replacing Electronics — it is growing the other three categories as a hedge. Growing non-Electronics from 12.8% to 20% would add ~$630K annual buffer.")
    return slide


def slide_10_regional_gaps(prs):
    """Slide 10 — Regional & Seasonal Gaps (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), RED_ALERT)
    section_label(slide, "Root Cause #3", "Regional & Seasonal Gaps", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "East is deteriorating fastest while South Wearables collapse goes unaddressed",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # Left — regional slope chart
    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(5.8), Inches(0.35),
                "Regional Growth Slopes (monthly trend)", font_size=12, bold=True, color=DARK_TEXT)

    reg_data = [
        ("West",  "+$609/mo",  609,   TEAL,      "-22.3% YoY"),
        ("North", "-$192/mo",  -192,  MID_GREY,  "-27.7% YoY"),
        ("South", "-$384/mo",  -384,  ORANGE,    "-27.8% YoY"),
        ("East",  "-$1,114/mo",-1114, RED_ALERT, "-35.3% YoY"),
    ]
    max_abs = 1114
    for i, (name, slope_lbl, slope, clr, yoy) in enumerate(reg_data):
        by = Inches(1.5) + i * Inches(0.95)
        add_textbox(slide, Inches(0.5), by + Inches(0.05),
                    Inches(1.2), Inches(0.3), name, font_size=11, bold=True, color=DARK_TEXT)
        # zero line at center
        center_x = Inches(1.85)
        max_bar = Inches(3.5)
        bar_w = max_bar * (abs(slope) / max_abs)
        if slope >= 0:
            add_rect(slide, center_x, by + Inches(0.1), bar_w, Inches(0.35), clr)
        else:
            add_rect(slide, center_x - bar_w, by + Inches(0.1), bar_w, Inches(0.35), clr)
        add_rect(slide, center_x - Inches(0.01), by, Inches(0.02), Inches(0.6), MID_GREY)
        add_textbox(slide, Inches(5.45), by + Inches(0.05),
                    Inches(1.5), Inches(0.3), slope_lbl, font_size=10, bold=True, color=clr)
        add_textbox(slide, Inches(5.45), by + Inches(0.38),
                    Inches(1.5), Inches(0.25), yoy, font_size=9, color=MID_GREY)

    # East callout
    add_rect(slide, Inches(0.5), Inches(5.4), Inches(5.9), Inches(0.75),
             RGBColor(255, 235, 235), line_color=RED_ALERT, line_width=1)
    add_textbox(slide, Inches(0.65), Inches(5.48), Inches(5.6), Inches(0.6),
                "East deterioration (-35.3% YoY) is nearly 2x the headline -28.4% — the aggregate figure conceals how severe East's collapse truly is.",
                font_size=10, color=RED_ALERT)

    # Right — South Wearables
    add_textbox(slide, Inches(7.0), Inches(1.05), Inches(5.8), Inches(0.35),
                "South Category Breakdown — Wearables Crisis", font_size=12, bold=True, color=DARK_TEXT)

    south_cat = [
        ("Wearables",   -58.6, -0.359, RED_ALERT, "PRIMARY DRIVER"),
        ("Audio",       -22.1, -0.308, ORANGE,    "Secondary drag"),
        ("Electronics",  -2.6, -0.145, MID_GREY,  "Essentially flat"),
        ("Accessories", +13.5, +0.124, TEAL,      "Growing — positive offset"),
    ]
    cat_h = Inches(0.9)
    for i, (cat, chg, corr, clr, note) in enumerate(south_cat):
        cy = Inches(1.5) + i * (cat_h + Inches(0.05))
        bg = RGBColor(255, 235, 235) if chg < -30 else (RGBColor(255, 245, 235) if chg < 0 else RGBColor(220, 248, 240))
        add_rect(slide, Inches(7.0), cy, Inches(5.8), cat_h, bg, line_color=clr, line_width=1)
        add_rect(slide, Inches(7.0), cy, Inches(0.06), cat_h, clr)
        add_textbox(slide, Inches(7.15), cy + Inches(0.08),
                    Inches(2.5), Inches(0.32), cat, font_size=12, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(7.15), cy + Inches(0.46),
                    Inches(2.5), Inches(0.28), note, font_size=9, color=MID_GREY)
        chg_str = f"{chg:+.1f}%"
        add_textbox(slide, Inches(11.0), cy + Inches(0.05),
                    Inches(1.7), Inches(0.42), chg_str, font_size=22, bold=True,
                    color=clr, align=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(11.0), cy + Inches(0.5),
                    Inches(1.7), Inches(0.28), f"corr={corr:+.3f}", font_size=8,
                    color=MID_GREY, align=PP_ALIGN.RIGHT)

    add_rect(slide, Inches(7.0), Inches(5.3), Inches(5.8), Inches(0.85),
             RGBColor(220, 248, 240), line_color=TEAL, line_width=1)
    add_textbox(slide, Inches(7.15), Inches(5.38), Inches(5.5), Inches(0.7),
                "SUPPLIER RULED OUT  All 5 suppliers have near-equal South revenue ($591K-$666K) and ratings (4.32-4.39). This is a demand-side problem, not a supply failure.",
                font_size=10, color=TEAL)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), RED_ALERT)
    add_notes(slide, "Two different problems requiring two different responses. East has a broad demand decline — it needs to understand what changed structurally and look at the West playbook. South Wearables is a discrete, identifiable signal: a specific category failing in a specific geography, with supplier quality explicitly ruled out. It is the kind of problem a regional investigation can solve in 90 days.")
    return slide


def slide_11_rejected(prs):
    """Slide 11 — Rejected Hypotheses (dark)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), RED_ALERT)
    section_label(slide, "Approach", "Hypotheses Tested", True)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Three hypotheses were rigorously tested and rejected with evidence",
                font_size=24, bold=True, color=WHITE)
    divider_line(slide, Inches(1.0), True)

    rejected = [
        {
            "id": "H-R1", "title": "NOT Supplier Quality",
            "accent": RED_ALERT, "verdict": "REJECTED",
            "hypothesis": "Supplier quality issues driving revenue volatility through elevated returns",
            "evidence": [
                "Return rates uniform: 2.5%-3.6% across all 5 suppliers",
                "Spread: only 1.04pp — systemic, not supplier-specific",
                "Feb 2022 & Apr 2023 troughs: no coincident return spikes",
                "Shipping cost: returned $7.76 vs non-returned $7.64 (negligible)",
            ],
            "verdict_note": "No correlation with revenue timing",
        },
        {
            "id": "H-R2", "title": "NOT Category Mix Shift",
            "accent": ORANGE, "verdict": "REJECTED",
            "hypothesis": "Category mix rotating from high-value Electronics to lower-value categories",
            "evidence": [
                "Electronics share range: 85.1%-88.7% (30 months)",
                "Standard deviation: only 0.9 ppts — frozen",
                "Mix effect contributes <1% of revenue variance",
                "No category rotation occurred at any point",
            ],
            "verdict_note": "Category shares structurally locked",
        },
        {
            "id": "H-R3", "title": "NOT Marketing Spend",
            "accent": PURPLE, "verdict": "REJECTED",
            "hypothesis": "Marketing spend changes drive revenue via lead-lag mechanism",
            "evidence": [
                "Marketing cost range: ~$4,600-$5,800/month (spread $1,200)",
                "Lag-0 r=0.62 (reactive co-movement)",
                "Lag-1 r=-0.33 (NEGATIVE — spend follows revenue)",
                "Google Ads $62,960 budget / ROI 32x vs Email 328x",
            ],
            "verdict_note": "Spend is reactive, not causal",
        },
    ]

    card_w = Inches(3.9)
    card_h = Inches(5.6)
    gap    = Inches(0.22)
    start_x = Inches(0.55)
    y0 = Inches(1.1)

    for i, c in enumerate(rejected):
        lx = start_x + i * (card_w + gap)
        add_rect(slide, lx, y0, card_w, card_h, RGBColor(14, 23, 45),
                 line_color=c["accent"], line_width=1)
        add_rect(slide, lx, y0, card_w, Inches(0.07), c["accent"])
        # Rejection X badge
        add_rect(slide, lx + card_w - Inches(0.6), y0 + Inches(0.1),
                 Inches(0.45), Inches(0.35), RED_ALERT)
        add_textbox(slide, lx + card_w - Inches(0.58), y0 + Inches(0.12),
                    Inches(0.4), Inches(0.3), "✗", font_size=14, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.18),
                    Inches(0.55), Inches(0.35), c["id"], font_size=13, bold=True, color=c["accent"])
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.58),
                    card_w - Inches(0.3), Inches(0.42), c["title"], font_size=15, bold=True, color=WHITE)
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(1.07),
                    card_w - Inches(0.3), Inches(0.42),
                    f"Hypothesis: {c['hypothesis']}", font_size=9, italic=True,
                    color=RGBColor(130, 150, 200))
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(1.52),
                    card_w - Inches(0.3), Inches(0.25),
                    "Evidence:", font_size=10, bold=True, color=RGBColor(200, 210, 240))
        for j, ev in enumerate(c["evidence"]):
            add_textbox(slide, lx + Inches(0.22), y0 + Inches(1.82) + j * Inches(0.67),
                        card_w - Inches(0.4), Inches(0.62),
                        f"• {ev}", font_size=10, color=RGBColor(175, 190, 225))
        add_rect(slide, lx, y0 + card_h - Inches(0.45), card_w, Inches(0.45), RGBColor(10, 18, 38))
        add_textbox(slide, lx + Inches(0.15), y0 + card_h - Inches(0.38),
                    card_w - Inches(0.25), Inches(0.32),
                    f"Rejected: {c['verdict_note']}", font_size=9, color=MID_GREY)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), RED_ALERT)
    add_notes(slide, "Ruling out hypotheses is analytically as important as confirming them — and for a business audience, it is even more important from a decision-making standpoint. These three rejected hypotheses represent actions that management might instinctively reach for. The data shows clearly that none of these would address the root cause. The marketing insight in particular deserves a moment: Google Ads is 10x less efficient per dollar than Email.")
    return slide


def slide_12_forecast(prs):
    """Slide 12 — Forecast (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    section_label(slide, "Predict", "Q4 2023 Forecast", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Q4 2023 forecast: $1.22M — decline continues without intervention",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # Forecast table
    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.35),
                "Monthly Forecast — Q4 2023 (Ensemble)", font_size=12, bold=True, color=DARK_TEXT)

    f_headers = ["Month", "Forecast", "Lower 95%", "Upper 95%"]
    f_rows = [
        ["Oct 2023", "$410,973", "$358,267", "$463,678"],
        ["Nov 2023", "$402,823", "$344,847", "$460,799"],
        ["Dec 2023", "$409,552", "$346,306", "$472,799"],
        ["Q4 TOTAL", "$1,223,348", "$1,049,420", "$1,397,276"],
    ]
    f_cw = [Inches(1.3), Inches(1.4), Inches(1.4), Inches(1.4)]
    f_cx = [Inches(0.5), Inches(1.85), Inches(3.3), Inches(4.75)]
    row_h = Inches(0.42)
    hy = Inches(1.48)
    add_rect(slide, Inches(0.5), hy, Inches(5.65), row_h, BLUE_PRIMARY)
    for j, (h, cw, cx) in enumerate(zip(f_headers, f_cw, f_cx)):
        add_textbox(slide, cx + Inches(0.05), hy + Inches(0.08),
                    cw - Inches(0.1), row_h - Inches(0.1),
                    h, font_size=10, bold=True, color=WHITE)
    for i, row in enumerate(f_rows):
        ry = hy + (i + 1) * row_h
        bg = RGBColor(220, 230, 255) if i == 3 else (RGBColor(228, 232, 248) if i % 2 == 0 else BG_LIGHT)
        add_rect(slide, Inches(0.5), ry, Inches(5.65), row_h, bg)
        for j, (cell, cw, cx) in enumerate(zip(row, f_cw, f_cx)):
            clr = BLUE_PRIMARY if i == 3 else DARK_TEXT
            bold = (i == 3) or (j == 1)
            add_textbox(slide, cx + Inches(0.05), ry + Inches(0.08),
                        cw - Inches(0.1), row_h - Inches(0.1),
                        cell, font_size=10, bold=bold, color=clr)

    # Model comparison
    add_textbox(slide, Inches(0.5), Inches(3.85), Inches(6.0), Inches(0.35),
                "Model Comparison — MAPE (lower = better)", font_size=12, bold=True, color=DARK_TEXT)

    model_data = [
        ("Linear Trend (Baseline)", 3.70, BLUE_BRIGHT,  "1,197,150"),
        ("ARIMA(1,0,0)",           3.92, PURPLE,        "1,241,707"),
        ("Holt-Winters",           4.00, TEAL,          "1,232,889"),
        ("Ensemble (Equal wt.)",   3.87, BLUE_PRIMARY,  "1,223,348"),
    ]
    max_mape = 4.5
    for i, (name, mape, clr, q4) in enumerate(model_data):
        my = Inches(4.3) + i * Inches(0.52)
        bar_w = Inches(4.5) * (mape / max_mape)
        add_rect(slide, Inches(0.5), my, Inches(4.5), Inches(0.42), RGBColor(228, 232, 248))
        add_rect(slide, Inches(0.5), my, bar_w, Inches(0.42), clr)
        add_textbox(slide, Inches(0.55), my + Inches(0.09),
                    bar_w - Inches(0.1), Inches(0.26), name, font_size=9,
                    color=WHITE if bar_w > Inches(1.5) else DARK_TEXT, bold=(i==3))
        add_textbox(slide, Inches(5.1), my + Inches(0.09),
                    Inches(0.6), Inches(0.26), f"{mape:.2f}%", font_size=10,
                    bold=True, color=clr)
        add_textbox(slide, Inches(5.8), my + Inches(0.09),
                    Inches(1.0), Inches(0.26), f"${q4}", font_size=9, color=MID_GREY)

    # Right — forecast visual and key messages
    add_textbox(slide, Inches(7.0), Inches(1.05), Inches(5.8), Inches(0.35),
                "Forecast Confidence & Key Messages", font_size=12, bold=True, color=DARK_TEXT)

    messages = [
        (TEAL,        "MAPE 3.7-4.0%", "All 3 models agree within 0.3 percentage points — strong directional consensus"),
        (BLUE_BRIGHT, "95% CI Width",  "Oct: $358K-$464K | Nov: $345K-$461K | Dec: $346K-$473K"),
        (RED_ALERT,   "No Reversal",   "Zero models detect a trend reversal — structural decline continues"),
        (AMBER,       "West Signal",   "+$609/month West growth — a region-level model would show stronger West"),
        (ORANGE,      "Data Limit",    "Only 21 monthly points — CI is intentionally wide; direction is unambiguous"),
    ]
    for i, (clr, lbl, txt) in enumerate(messages):
        my = Inches(1.5) + i * Inches(1.0)
        add_rect(slide, Inches(7.0), my, Inches(5.8), Inches(0.9), RGBColor(228, 232, 248),
                 line_color=clr, line_width=1)
        add_rect(slide, Inches(7.0), my, Inches(0.06), Inches(0.9), clr)
        add_textbox(slide, Inches(7.15), my + Inches(0.06),
                    Inches(5.5), Inches(0.3), lbl, font_size=11, bold=True, color=clr)
        add_textbox(slide, Inches(7.15), my + Inches(0.42),
                    Inches(5.5), Inches(0.38), txt, font_size=10, color=DARK_TEXT)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    add_notes(slide, "The forecast is a confirmation, not a surprise — it follows directly from the structural decline already established in the data. The value to the audience is not the point estimate but the model consensus: three independent modelling approaches all land within 0.3 MAPE points of each other, and all point the same direction. Q4 will be lower than Q3 unless something changes.")
    return slide


def slide_13_conclusions(prs):
    """Slide 13 — Conclusions (dark)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    section_label(slide, "Conclusions", "Key Takeouts", True)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Three conclusions that demand immediate management action",
                font_size=24, bold=True, color=WHITE)
    divider_line(slide, Inches(1.0), True)

    conclusions = [
        {
            "num": "01", "title": "Act Now — Volume",
            "accent": RED_ALERT, "urgency": "P1 — Highest Urgency",
            "bullets": [
                "Volume drives 70% of variance — confirmed, r=0.818",
                "Email 328x ROI vs Google Ads 32x — reallocate $20-25K",
                "Activate Feb North promotions before Q1 2024",
                "Protect 1,950 high-value orders ($1,984 AOV)",
            ],
            "footer": "Revenue impact in 60-90 days · no capital required",
        },
        {
            "num": "02", "title": "Invest — Diversification",
            "accent": AMBER, "urgency": "P2 — Medium-Term",
            "bullets": [
                "87% Electronics = single point of failure",
                "Accessories 40.7% margin / Audio 38.3% / Wearables 37.9%",
                "Grow non-Electronics from 12.8% to 20% share",
                "+$630K annual buffer when threshold reached",
            ],
            "footer": "6-12 month portfolio reshaping · board-level decision",
        },
        {
            "num": "03", "title": "Investigate — South",
            "accent": TEAL, "urgency": "P3 — Near-Term",
            "bullets": [
                "South Wearables -58.6%, South Audio -22.1%",
                "Supplier quality ruled out — all 5 equal in South",
                "Commission demand investigation in 30 days",
                "West +$609/month is the replication model",
            ],
            "footer": "Investigate within 30 days · 90-day resolution target",
        },
    ]

    card_w = Inches(3.9)
    card_h = Inches(5.6)
    gap    = Inches(0.22)
    start_x = Inches(0.55)
    y0 = Inches(1.1)

    for i, c in enumerate(conclusions):
        lx = start_x + i * (card_w + gap)
        add_rect(slide, lx, y0, card_w, card_h, RGBColor(14, 23, 45),
                 line_color=c["accent"], line_width=1)
        add_rect(slide, lx, y0, card_w, Inches(0.07), c["accent"])
        # Urgency badge
        add_rect(slide, lx + card_w - Inches(1.45), y0 + Inches(0.12),
                 Inches(1.3), Inches(0.3), c["accent"])
        add_textbox(slide, lx + card_w - Inches(1.43), y0 + Inches(0.14),
                    Inches(1.26), Inches(0.26), c["urgency"][:2], font_size=10,
                    bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.18),
                    Inches(0.55), Inches(0.35), c["num"], font_size=20, bold=True, color=c["accent"])
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.58),
                    card_w - Inches(0.3), Inches(0.42), c["title"], font_size=16, bold=True, color=WHITE)
        for j, bullet in enumerate(c["bullets"]):
            add_textbox(slide, lx + Inches(0.22), y0 + Inches(1.15) + j * Inches(0.85),
                        card_w - Inches(0.4), Inches(0.78),
                        f"• {bullet}", font_size=11, color=RGBColor(180, 195, 230))
        add_rect(slide, lx, y0 + card_h - Inches(0.52), card_w, Inches(0.52), RGBColor(10, 18, 38))
        add_textbox(slide, lx + Inches(0.15), y0 + card_h - Inches(0.44),
                    card_w - Inches(0.25), Inches(0.38), c["footer"], font_size=9, color=MID_GREY)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), BLUE_PRIMARY)
    add_notes(slide, "Three levels of urgency. The volume and marketing reallocation finding is actionable this quarter with no capital investment — it is purely a budget allocation decision. The diversification investment requires strategic planning. The South investigation needs to be commissioned immediately even if it takes 90 days to complete, because the Wearables collapse is ongoing.")
    return slide


def slide_14_roadmap(prs):
    """Slide 14 — Strategy Roadmap (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), TEAL)
    section_label(slide, "Action", "Strategy Maps", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "A three-phase roadmap from immediate triage to structural resilience",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    phases = [
        {
            "phase": "Phase 1", "title": "Triage",
            "timeline": "0-90 Days", "accent": RED_ALERT,
            "actions": [
                "Shift $20-25K Google Ads → Email & Referral",
                "Design Feb North promotion campaign",
                "Identify & protect 1,950 high-value orders",
                "Set monthly forecast tracking baseline ($1.22M Q4)",
            ],
            "outcome": "Revenue & margin recovery in current quarter",
        },
        {
            "phase": "Phase 2", "title": "Investigate",
            "timeline": "90-180 Days", "accent": AMBER,
            "actions": [
                "Commission South Wearables demand investigation",
                "Extract West growth playbook (+$609/month)",
                "Pilot West tactics in North (-$192/month slope)",
                "Flag any month >15% miss for immediate review",
            ],
            "outcome": "Understand regional and category demand barriers",
        },
        {
            "phase": "Phase 3", "title": "Structural",
            "timeline": "6-18 Months", "accent": TEAL,
            "actions": [
                "Grow non-Electronics from 12.8% to 20% share",
                "Build region-level forecast models (4 regions)",
                "Reduce revenue CV from 6.5% toward < 5%",
                "Improve margin mix from 33.2% toward 34-35%",
            ],
            "outcome": "Portfolio resilience and forecast accuracy",
        },
    ]

    card_w = Inches(3.9)
    card_h = Inches(5.55)
    gap    = Inches(0.22)
    start_x = Inches(0.55)
    y0 = Inches(1.1)

    for i, p in enumerate(phases):
        lx = start_x + i * (card_w + gap)
        add_rect(slide, lx, y0, card_w, card_h, RGBColor(235, 238, 252),
                 line_color=p["accent"], line_width=2)
        # Phase banner
        add_rect(slide, lx, y0, card_w, Inches(0.65), p["accent"])
        add_textbox(slide, lx + Inches(0.15), y0 + Inches(0.04),
                    card_w * 0.5 - Inches(0.1), Inches(0.3),
                    p["phase"], font_size=12, bold=True, color=WHITE)
        add_textbox(slide, lx + Inches(0.15), y0 + Inches(0.34),
                    card_w - Inches(0.3), Inches(0.28),
                    p["timeline"], font_size=11, color=RGBColor(220, 230, 255))
        add_textbox(slide, lx + Inches(0.18), y0 + Inches(0.78),
                    card_w - Inches(0.3), Inches(0.42),
                    p["title"], font_size=20, bold=True, color=p["accent"])
        for j, action in enumerate(p["actions"]):
            add_textbox(slide, lx + Inches(0.2), y0 + Inches(1.3) + j * Inches(0.85),
                        card_w - Inches(0.35), Inches(0.78),
                        f"• {action}", font_size=11, color=DARK_TEXT)
        add_rect(slide, lx, y0 + card_h - Inches(0.52), card_w, Inches(0.52),
                 RGBColor(215, 220, 245))
        add_textbox(slide, lx + Inches(0.15), y0 + card_h - Inches(0.44),
                    card_w - Inches(0.25), Inches(0.38),
                    p["outcome"], font_size=9, color=MID_GREY)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), TEAL)
    add_notes(slide, "Phase 1 is actionable in the next 90 days with no capital investment — it is purely reallocation decisions and promotional planning. Phase 2 requires time and research investment but can be commissioned immediately. Phase 3 requires strategic budget and board-level decision, but the groundwork starts now with understanding what West is doing right.")
    return slide


def slide_15_recommendations(prs):
    """Slide 15 — Recommendations Table (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), TEAL)
    section_label(slide, "Action", "Recommendation", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Five prioritised actions ranked by revenue impact and effort",
                font_size=20, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # Priority legend
    legend = [
        ("P1-Now",    RED_ALERT, "Immediate — no capital"),
        ("P2-Soon",   AMBER,     "Near-term — research"),
        ("P3-Invest", TEAL,      "Strategic — 12 months"),
    ]
    lx = Inches(0.5)
    for lbl, clr, note in legend:
        add_rect(slide, lx, Inches(1.08), Inches(0.8), Inches(0.3), clr)
        add_textbox(slide, lx + Inches(0.04), Inches(1.1), Inches(0.72), Inches(0.25),
                    lbl, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.85), Inches(1.1), Inches(1.6), Inches(0.25),
                    note, font_size=8, color=MID_GREY)
        lx += Inches(2.6)

    # Table headers
    col_defs = [
        ("Priority", Inches(0.85)),
        ("Action", Inches(4.2)),
        ("Owner", Inches(1.5)),
        ("Timeline", Inches(0.95)),
        ("Expected Impact", Inches(2.7)),
        ("Success Metric", Inches(2.6)),
    ]
    col_starts = []
    cx = Inches(0.25)
    for _, w in col_defs:
        col_starts.append(cx)
        cx += w + Inches(0.05)

    row_h = Inches(0.42)
    hy = Inches(1.45)
    add_rect(slide, Inches(0.25), hy, Inches(12.8), row_h, DARK_TEXT)
    for j, ((hdr, w), sx) in enumerate(zip(col_defs, col_starts)):
        add_textbox(slide, sx + Inches(0.04), hy + Inches(0.08),
                    w - Inches(0.08), row_h - Inches(0.12),
                    hdr, font_size=9, bold=True, color=WHITE)

    rows = [
        ("P1-Now",    "Reallocate $20-25K from Google Ads to Email & Referral channels",
         "Marketing", "30 days", "+40-60% blended paid ROI", "Paid ROI > 80x blended",
         RED_ALERT),
        ("P1-Now",    "Build February North region promotion to buffer structural -36.9% trough",
         "Marketing", "60 days", "Recover 30-50% of Feb shortfall", "Feb MoM drop < -10%",
         RED_ALERT),
        ("P2-Soon",   "Commission South Wearables demand investigation (regional & competitor)",
         "Product/Regional", "90 days", "Recover $50-100K/year Wearables South", "South trend reversal",
         AMBER),
        ("P2-Soon",   "Extract West growth playbook and pilot in North region",
         "Regional Mgr", "180 days", "North slope improvement +$300/month", "North slope turns positive",
         AMBER),
        ("P3-Invest", "Invest in Accessories & Audio to grow non-Electronics from 13% to 20%",
         "Product/Strategy", "12 months", "Reduce revenue CV, improve margin mix", "Non-Electronics > 18%",
         TEAL),
    ]

    for i, (pri, action, owner, timeline, impact, metric, clr) in enumerate(rows):
        ry = hy + (i + 1) * row_h
        bg = RGBColor(255, 235, 235) if clr == RED_ALERT else (
             RGBColor(255, 248, 230) if clr == AMBER else RGBColor(225, 248, 242))
        add_rect(slide, Inches(0.25), ry, Inches(12.8), row_h, bg)
        # Priority badge
        add_rect(slide, Inches(0.25), ry, Inches(0.85), row_h, clr)
        add_textbox(slide, Inches(0.27), ry + Inches(0.06),
                    Inches(0.78), row_h - Inches(0.1),
                    pri, font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cells = [None, action, owner, timeline, impact, metric]
        for j, (cell, (_, w), sx) in enumerate(zip(cells, col_defs, col_starts)):
            if cell is None:
                continue
            color = DARK_TEXT
            add_textbox(slide, sx + Inches(0.04), ry + Inches(0.06),
                        w - Inches(0.08), row_h - Inches(0.1),
                        cell, font_size=9, color=color)

    # Footer
    add_textbox(slide, Inches(0.25), Inches(6.65), Inches(12.8), Inches(0.35),
                "P1 items require zero capital — budget reallocation only. P2 can be scoped and commissioned immediately. P3 requires strategic investment over 12 months.",
                font_size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), TEAL)
    add_notes(slide, "P1 items require no capital, just reallocation decisions — these can be decided and assigned today. P2 items require time and research but can be scoped and commissioned now. P3 requires strategic budget and product investment over a 12-month horizon. The table is designed for the management team to assign owners and timelines before leaving this meeting.")
    return slide


def slide_16_impact(prs):
    """Slide 16 — Impact (light)"""
    slide = blank_slide(prs)
    set_bg(slide, BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), TEAL)
    section_label(slide, "Action", "Expected Impact", False)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.2), Inches(0.5),
                "Executing the roadmap converts a declining trend into a stabilised and growing business",
                font_size=19, bold=True, color=DARK_TEXT)
    divider_line(slide, Inches(0.95), False)

    # Left — Short-term recovery
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(5.9), Inches(5.45),
             RGBColor(220, 228, 255), line_color=BLUE_PRIMARY, line_width=2)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(5.9), Inches(0.5), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.65), Inches(1.1), Inches(5.6), Inches(0.38),
                "Short-Term Revenue Recovery (0-90 days)", font_size=13, bold=True, color=WHITE)

    st_items = [
        ("Q4 Base Case",      "$1,223,348",   MID_GREY,  "No action scenario"),
        ("Marketing ROI gain","+40-60%",       GREEN_OK,  "Shift $20-25K to Email/Referral"),
        ("Feb trough recovery","30-50%",        TEAL,      "North promotion campaign"),
        ("High-value protect", "1,950 orders",  AMBER,     "= 975 standard orders in revenue"),
        ("Margin recovery",   "$15-25K/year",  BLUE_BRIGHT,"Annual margin uplift from reallocation"),
    ]
    for i, (lbl, val, clr, note) in enumerate(st_items):
        iy = Inches(1.68) + i * Inches(0.82)
        add_rect(slide, Inches(0.6), iy, Inches(5.7), Inches(0.72), RGBColor(235, 240, 255))
        add_rect(slide, Inches(0.6), iy, Inches(0.05), Inches(0.72), clr)
        add_textbox(slide, Inches(0.72), iy + Inches(0.06),
                    Inches(3.0), Inches(0.28), lbl, font_size=10, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(0.72), iy + Inches(0.37),
                    Inches(3.0), Inches(0.25), note, font_size=9, color=MID_GREY)
        add_textbox(slide, Inches(3.85), iy + Inches(0.08),
                    Inches(2.3), Inches(0.45), val, font_size=16, bold=True,
                    color=clr, align=PP_ALIGN.RIGHT)

    # Right — Structural resilience
    add_rect(slide, Inches(6.9), Inches(1.05), Inches(5.9), Inches(5.45),
             RGBColor(215, 248, 242), line_color=TEAL, line_width=2)
    add_rect(slide, Inches(6.9), Inches(1.05), Inches(5.9), Inches(0.5), TEAL)
    add_textbox(slide, Inches(7.05), Inches(1.1), Inches(5.6), Inches(0.38),
                "Structural Resilience (6-18 months)", font_size=13, bold=True, color=WHITE)

    lt_items = [
        ("Base Case Q4 2024",     "$1.3-1.4M",   TEAL,       "With P1+P2 executed"),
        ("Revenue CV target",     "< 5%",         BLUE_BRIGHT,"Down from 6.5% — reduced volatility"),
        ("Margin mix target",     "34-35%",       GREEN_OK,   "Up from 33.2% — diversification benefit"),
        ("Non-Electronics share", "> 18%",        AMBER,      "12-month target, up from 12.8%"),
        ("Early warning trigger", "> 15% miss",   RED_ALERT,  "Immediate diagnostic review required"),
    ]
    for i, (lbl, val, clr, note) in enumerate(lt_items):
        iy = Inches(1.68) + i * Inches(0.82)
        add_rect(slide, Inches(7.0), iy, Inches(5.7), Inches(0.72), RGBColor(225, 252, 246))
        add_rect(slide, Inches(7.0), iy, Inches(0.05), Inches(0.72), clr)
        add_textbox(slide, Inches(7.12), iy + Inches(0.06),
                    Inches(3.0), Inches(0.28), lbl, font_size=10, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(7.12), iy + Inches(0.37),
                    Inches(3.0), Inches(0.25), note, font_size=9, color=MID_GREY)
        add_textbox(slide, Inches(10.25), iy + Inches(0.08),
                    Inches(2.3), Inches(0.45), val, font_size=16, bold=True,
                    color=clr, align=PP_ALIGN.RIGHT)

    add_rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), TEAL)
    add_notes(slide, "Close the presentation by bringing the audience back to the business stakes. The data analysis is not academic — the Q4 forecast shows a declining trend that will compound if unaddressed. The roadmap is not hypothetical — it is directly connected to the diagnosed root causes. The margin staying at 33.2% throughout is actually a strategic asset — it tells management the business model is sound and the problem is purely about demand generation and portfolio resilience. Both are solvable.")
    return slide


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    print("Building slide 01 — Cover...")
    slide_01_cover(prs)
    print("Building slide 02 — Background...")
    slide_02_background(prs)
    print("Building slide 03 — Framework...")
    slide_03_framework(prs)
    print("Building slide 04 — YoY Comparison...")
    slide_04_yoy(prs)
    print("Building slide 05 — Category & Region...")
    slide_05_category_region(prs)
    print("Building slide 06 — Segment & Seasonality...")
    slide_06_segment(prs)
    print("Building slide 07 — Findings Overview...")
    slide_07_findings_overview(prs)
    print("Building slide 08 — Volume Decomposition...")
    slide_08_volume_decomp(prs)
    print("Building slide 09 — Concentration Risk...")
    slide_09_concentration(prs)
    print("Building slide 10 — Regional & Seasonal Gaps...")
    slide_10_regional_gaps(prs)
    print("Building slide 11 — Rejected Hypotheses...")
    slide_11_rejected(prs)
    print("Building slide 12 — Forecast...")
    slide_12_forecast(prs)
    print("Building slide 13 — Conclusions...")
    slide_13_conclusions(prs)
    print("Building slide 14 — Strategy Roadmap...")
    slide_14_roadmap(prs)
    print("Building slide 15 — Recommendations Table...")
    slide_15_recommendations(prs)
    print("Building slide 16 — Impact...")
    slide_16_impact(prs)

    out_path = r"c:\This PC\the future analyst\AIDA\AIDA-All\AIDA\ai_analyst\data\reports\revenue\techworld_data\slidedeck.pptx"
    prs.save(out_path)
    print(f"\nSaved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")
    return out_path


if __name__ == "__main__":
    build_deck()
